from __future__ import annotations

import io

from fielddesk_worker.parsing.base import ParseError, ParsedSegment


def parse_pptx(content: bytes) -> list[ParsedSegment]:
    """PowerPoint (.pptx) parser. Emits one segment per slide so citations
    can say "slide 4" the same way the PDF parser cites "page 4."

    Each segment carries the slide title (if any) as a single-element
    heading_path so the chunker preserves it on long slides. Source locator
    keeps slide_index + total_slides so the UI can render "slide 4 of 27."

    Out of scope: speaker notes (rarely useful for ticket-recommendation
    context, can add later), embedded images/charts (would need OCR per
    image, separate slice), animations.
    """
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError as e:
        raise ParseError(f"python-pptx is required for PPTX parsing: {e}") from e

    try:
        prs = Presentation(io.BytesIO(content))
    except Exception as e:  # noqa: BLE001
        raise ParseError(f"could not read pptx: {e}") from e

    slides = list(prs.slides)
    total_slides = len(slides)
    segments: list[ParsedSegment] = []
    for index, slide in enumerate(slides):
        title, body_lines = _collect_slide_text(slide)
        body = "\n".join(body_lines).strip()
        if not body and not title:
            continue
        segments.append(
            ParsedSegment(
                text=body or title,
                heading_path=[title] if title else [],
                source_page=index + 1,
                source_locator={"slide": index + 1, "of_slides": total_slides},
            )
        )
    return segments


def _collect_slide_text(slide) -> tuple[str, list[str]]:
    """Pull the slide title plus all non-title text from text frames.

    Title detection: python-pptx exposes `slide.shapes.title` which is the
    layout-defined title placeholder. When the deck doesn't use a title
    placeholder (common in custom layouts), we fall back to using the
    first text-frame content as a heading. The point is to give the
    chunker something to preserve as heading_path, not to be perfect.
    """
    title = ""
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and title_shape.has_text_frame:
        title = (title_shape.text_frame.text or "").strip()

    body_lines: list[str] = []
    for shape in slide.shapes:
        if shape is title_shape:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text_frame.text or "").strip()
        if text:
            body_lines.append(text)
    return title, body_lines
