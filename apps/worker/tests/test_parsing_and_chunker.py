from __future__ import annotations

import sys
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

try:
    from fielddesk_worker.embeddings.chunker import chunk_segments
    from fielddesk_worker.parsing import (
        SUPPORTED_MIME_TYPES,
        ParsedSegment,
        parse_document,
    )
    from fielddesk_worker.parsing.markdown import parse_markdown
    from fielddesk_worker.parsing.text import parse_text
except ModuleNotFoundError as exc:
    raise unittest.SkipTest(f"worker dependencies are not installed: {exc.name}") from exc

try:
    import tiktoken  # noqa: F401

    _HAS_TIKTOKEN = True
except ModuleNotFoundError:
    _HAS_TIKTOKEN = False


class TextParserTests(unittest.TestCase):
    def test_strips_whitespace_and_returns_one_segment(self):
        segments = parse_text(b"   hello world   \n\n")
        self.assertEqual(len(segments), 1)
        self.assertEqual(segments[0].text, "hello world")
        self.assertEqual(segments[0].heading_path, [])

    def test_empty_returns_no_segments(self):
        self.assertEqual(parse_text(b""), [])
        self.assertEqual(parse_text(b"\n   \n"), [])

    def test_falls_back_to_latin1_on_bad_utf8(self):
        # 0xa1 is invalid utf-8 but valid latin-1 (¡).
        segments = parse_text(b"\xa1hola")
        self.assertEqual(len(segments), 1)


class MarkdownParserTests(unittest.TestCase):
    def test_one_heading_one_section(self):
        md = b"# Title\n\nBody text here.\n"
        segments = parse_markdown(md)
        self.assertEqual(len(segments), 1)
        self.assertEqual(segments[0].heading_path, ["Title"])
        self.assertIn("Body text", segments[0].text)

    def test_nested_headings_carry_path(self):
        md = (
            b"# Hydraulic Pumps\n"
            b"intro paragraph\n"
            b"## Troubleshooting\n"
            b"### Pressure Loss\n"
            b"Step 1: check the inlet.\n"
        )
        segments = parse_markdown(md)
        # 3 sections: under "Hydraulic Pumps", under "+Troubleshooting"
        # (empty body so no segment), under "+Pressure Loss". The empty
        # heading should not emit an empty-body segment.
        self.assertGreaterEqual(len(segments), 2)
        last = segments[-1]
        self.assertEqual(last.heading_path, ["Hydraulic Pumps", "Troubleshooting", "Pressure Loss"])

    def test_headings_inside_fenced_blocks_are_ignored(self):
        md = (
            b"# Real Heading\n"
            b"```\n"
            b"# This looks like a heading but isn't\n"
            b"```\n"
            b"after fence\n"
        )
        segments = parse_markdown(md)
        # All text belongs under "Real Heading"; the fence inside should
        # NOT have created a second section.
        self.assertEqual(len(segments), 1)
        self.assertEqual(segments[0].heading_path, ["Real Heading"])

    def test_dedented_heading_truncates_stack(self):
        # Going from H3 back up to H2 should pop the H3 off the stack.
        md = (
            b"## Section A\n"
            b"### Sub A.1\n"
            b"first body\n"
            b"## Section B\n"
            b"second body\n"
        )
        segments = parse_markdown(md)
        paths = [s.heading_path for s in segments]
        self.assertIn(["Section A", "Sub A.1"], paths)
        self.assertIn(["Section B"], paths)


@unittest.skipUnless(_HAS_TIKTOKEN, "tiktoken is not installed")
class ChunkerTests(unittest.TestCase):
    def test_short_segment_passes_through_unchanged(self):
        segments = [
            ParsedSegment(text="quick brown fox", heading_path=["A"], source_page=1)
        ]
        chunks = chunk_segments(segments)
        self.assertEqual(len(chunks), 1)
        self.assertEqual(chunks[0].text, "quick brown fox")
        self.assertEqual(chunks[0].heading_path, ["A"])
        self.assertEqual(chunks[0].source_page, 1)
        # chunk_index is monotonically increasing from 0.
        self.assertEqual(chunks[0].chunk_index, 0)

    def test_long_segment_splits_with_overlap(self):
        # Build a body well above 512 tokens by repeating a sentence.
        sentence = "This is a test sentence for the chunker. " * 200
        segments = [ParsedSegment(text=sentence, heading_path=["Big"])]
        chunks = chunk_segments(segments, target_tokens=512, overlap_tokens=64)
        self.assertGreater(len(chunks), 1)
        # Each chunk should respect the token budget (allow some slack for
        # the recursive splitter not landing exactly on boundary).
        for c in chunks:
            self.assertLessEqual(c.token_count, 600)
            self.assertEqual(c.heading_path, ["Big"])
        # Indexes are sequential.
        self.assertEqual([c.chunk_index for c in chunks], list(range(len(chunks))))

    def test_content_hash_is_stable_for_same_input(self):
        s = ParsedSegment(text="hello world", heading_path=["A"], source_page=2)
        c1 = chunk_segments([s])
        c2 = chunk_segments([s])
        self.assertEqual(c1[0].content_hash, c2[0].content_hash)

    def test_content_hash_differs_when_page_differs(self):
        # Same text on two different pages must produce distinct hashes,
        # otherwise the idempotent UPSERT would silently drop a chunk.
        a = ParsedSegment(text="boilerplate", source_page=1)
        b = ParsedSegment(text="boilerplate", source_page=2)
        ca = chunk_segments([a])[0]
        cb = chunk_segments([b])[0]
        self.assertNotEqual(ca.content_hash, cb.content_hash)

    def test_empty_segments_produce_zero_chunks(self):
        self.assertEqual(chunk_segments([]), [])
        self.assertEqual(chunk_segments([ParsedSegment(text="")]), [])

    def test_rejects_bad_params(self):
        with self.assertRaises(ValueError):
            chunk_segments([ParsedSegment(text="x")], target_tokens=0)
        with self.assertRaises(ValueError):
            chunk_segments([ParsedSegment(text="x")], target_tokens=10, overlap_tokens=10)


class PdfParserTests(unittest.TestCase):
    def test_blank_scanned_pdf_raises_after_ocr_returns_empty(self):
        # Blank pages go through the OCR fallback (which now exists) and
        # tesseract returns the empty string. parse_pdf must surface this
        # as ParseError so the document lands in status='failed' with a
        # parse_error explaining that OCR ran but found nothing —
        # critically not as a silently-empty 'ready' document.
        #
        # The OCR fallback's import line raises ParseError if pypdfium2
        # or pytesseract are absent; either path lands here, which is
        # exactly what we want — a clear failure rather than a silent
        # success.
        try:
            from pypdf import PdfWriter
        except ModuleNotFoundError:
            self.skipTest("pypdf not installed")
        from fielddesk_worker.parsing import ParseError
        from fielddesk_worker.parsing.pdf import parse_pdf
        import io as _io

        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        writer.add_blank_page(width=612, height=792)
        buf = _io.BytesIO()
        writer.write(buf)

        with self.assertRaises(ParseError) as ctx:
            parse_pdf(buf.getvalue())
        msg = str(ctx.exception).lower()
        # Either "OCR returned no text" (deps installed, blank pages) or
        # "OCR fallback requires" (deps missing in local dev) — both are
        # legitimate failure paths; what we're guarding is "must not
        # silently succeed."
        self.assertTrue(
            "ocr" in msg or "no extractable text" in msg,
            f"expected OCR-related ParseError, got: {ctx.exception}",
        )

    def test_encrypted_pdf_message_is_actionable(self):
        # The error_message ends up in documents.parse_error and is
        # surfaced verbatim in the failures dashboard. Operators reading
        # "encrypted PDFs are not supported in v1" don't know what to do;
        # the message must tell them how to fix it.
        try:
            from pypdf import PdfWriter
        except ModuleNotFoundError:
            self.skipTest("pypdf not installed")
        from fielddesk_worker.parsing import ParseError
        from fielddesk_worker.parsing.pdf import parse_pdf
        import io as _io

        writer = PdfWriter()
        writer.add_blank_page(width=612, height=792)
        writer.encrypt(user_password="secret")
        buf = _io.BytesIO()
        writer.write(buf)

        with self.assertRaises(ParseError) as ctx:
            parse_pdf(buf.getvalue())
        msg = str(ctx.exception).lower()
        self.assertIn("password", msg)
        # The actionable phrasing: tell the operator what to do.
        self.assertTrue(
            "remove" in msg or "unencrypted" in msg,
            f"encrypted-PDF error must tell the operator how to fix it, got: {ctx.exception}",
        )


class PptxParserTests(unittest.TestCase):
    def test_emits_one_segment_per_non_empty_slide(self):
        try:
            from pptx import Presentation
        except ModuleNotFoundError:
            self.skipTest("python-pptx not installed")
        from fielddesk_worker.parsing.pptx import parse_pptx
        import io as _io

        prs = Presentation()
        # Slide 1: title + body.
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Hydraulic Pump Maintenance"
        slide.placeholders[1].text = "Check the inlet filter every 90 days."
        # Slide 2: title only.
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2.shapes.title.text = "Safety Notes"
        buf = _io.BytesIO()
        prs.save(buf)

        segments = parse_pptx(buf.getvalue())
        self.assertEqual(len(segments), 2)
        # First slide carries the title as heading_path and includes the
        # body content.
        self.assertEqual(segments[0].heading_path, ["Hydraulic Pump Maintenance"])
        self.assertIn("inlet filter", segments[0].text)
        # source_locator carries slide numbering so citations can render
        # "slide N of M".
        self.assertEqual(segments[0].source_locator["slide"], 1)
        self.assertEqual(segments[0].source_locator["of_slides"], 2)
        # Second slide (title only) still emits a segment so the title
        # is searchable even without body content.
        self.assertEqual(segments[1].heading_path, ["Safety Notes"])


class DocParserTests(unittest.TestCase):
    def test_missing_libreoffice_raises_actionable_error(self):
        # Same posture as the encrypted-PDF message: if soffice is absent,
        # the failure must name the missing dep so an operator can fix
        # the deploy rather than guessing what "could not parse .doc"
        # means.
        try:
            from fielddesk_worker.parsing import ParseError
            from fielddesk_worker.parsing.doc import parse_doc
        except ModuleNotFoundError:
            self.skipTest("worker deps not installed")
        import shutil
        from unittest.mock import patch

        with patch.object(shutil, "which", return_value=None):
            with self.assertRaises(ParseError) as ctx:
                parse_doc(b"\xd0\xcf\x11\xe0fake-ole-header")
            msg = str(ctx.exception).lower()
            self.assertIn("libreoffice", msg)


class ParseDocumentRouterTests(unittest.TestCase):
    def test_supported_mime_types_match_what_router_knows(self):
        # If this set changes, the Go handler's allowedDocumentMimes AND
        # the frontend's ACCEPT_HINT / EXT_TO_MIME need the same update —
        # the test exists to remind us. Three sync points, no shared schema.
        self.assertEqual(
            set(SUPPORTED_MIME_TYPES.keys()),
            {
                "text/plain",
                "text/markdown",
                "text/x-markdown",
                "application/pdf",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/msword",
            },
        )

    def test_unsupported_mime_raises(self):
        from fielddesk_worker.parsing import ParseError

        with self.assertRaises(ParseError):
            parse_document(b"<x/>", "application/xml")

    def test_routes_text_correctly(self):
        segments = parse_document(b"hello", "text/plain")
        self.assertEqual(len(segments), 1)
        self.assertEqual(segments[0].text, "hello")


if __name__ == "__main__":
    unittest.main()
